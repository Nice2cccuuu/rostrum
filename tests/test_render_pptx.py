"""Tests for glyph measurement, template ingestion, binding and rendering.

The load-bearing assertion in this file is that **measured capacity is a sound
upper bound**: text filled to exactly a slot's reported capacity must fit. The
planner treats capacity as a hard constraint, so a capacity that overflows is
worse than no capacity at all.
"""

from __future__ import annotations

from typing import ClassVar

import pytest

from rostrum.budget import allocate
from rostrum.ir import (
    Asset,
    AssetKind,
    Block,
    BlockType,
    Channel,
    Deck,
    DeckMeta,
    DeliveryPlan,
    Density,
    Section,
    Slide,
    SlideRole,
    SourceDocument,
    SourceSpan,
)
from rostrum.measure.text import (
    FontMetrics,
    _filler_of,
    capacity_units,
    emu_to_pt,
    is_ideograph,
    lines_available,
    load_font,
    measure_text,
    pt_to_emu,
    wrap_text,
)
from rostrum.templates import bind, capacity_caps, ingest_pptx, overflow_rate

pptx = pytest.importorskip("pptx", reason="python-pptx not installed")

CJK_FONT = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
DOC = "src"


@pytest.fixture(scope="module")
def font() -> FontMetrics:
    import os

    if not os.path.exists(CJK_FONT):
        pytest.skip("CJK font unavailable")
    return load_font(CJK_FONT)


@pytest.fixture(scope="module")
def template(tmp_path_factory):
    """The stock PowerPoint template, ingested once."""
    from pptx import Presentation

    path = tmp_path_factory.mktemp("tpl") / "default.pptx"
    Presentation().save(str(path))
    contract, report = ingest_pptx(
        str(path), name="Default", font_path=CJK_FONT, language="zh"
    )
    return contract, report


@pytest.fixture
def figure(tmp_path):
    from PIL import Image

    p = tmp_path / "fig.png"
    Image.new("RGB", (1200, 750), "white").save(p)
    return str(p)


def bullet(text: str, importance: float = 0.6, **kw) -> Block:
    return Block(
        type=BlockType.BULLET,
        content=text,
        importance=importance,
        spans=[SourceSpan(doc_id=DOC, start=0, end=20)],
        **kw,
    )


def make_deck(slides: list[Slide], *, seconds: int = 480, **plan) -> Deck:
    return Deck(
        meta=DeckMeta(title="低资源表征学习", language="zh"),
        delivery=DeliveryPlan(total_seconds=seconds, **plan),
        sources=[SourceDocument(doc_id=DOC, char_count=50_000)],
        sections=[Section(title="S", slides=slides)],
    )


# --------------------------------------------------------------------------- #
# Glyph metrics
# --------------------------------------------------------------------------- #


class TestGlyphMetrics:
    def test_real_advances_refute_the_common_heuristic(self, font):
        """"A CJK char is two Latin chars" is wrong in both directions.

        Ideographs are ~1.0 em, while Latin advances vary threefold between
        ``i`` and ``W``. No single ratio can stand in for measurement.
        """
        assert font.measured
        assert font.advance_em("研") == pytest.approx(1.0, abs=0.02)
        latin = [font.advance_em(c) for c in "AiW"]
        assert max(latin) / min(latin) > 2.5
        ratio = font.advance_em("研") / (sum(latin) / len(latin))
        assert 1.5 < ratio < 1.9  # not 2.0

    def test_fallback_is_usable_without_a_font_file(self):
        """Accuracy degrades; nothing breaks."""
        f = FontMetrics(None)
        assert not f.measured
        assert f.advance_em("研") == pytest.approx(1.0)
        assert 0 < f.advance_em("a") < 1.0

    def test_unmapped_codepoint_falls_back(self, font):
        assert font.advance_em("\ue000") > 0

    def test_emu_point_roundtrip(self):
        assert emu_to_pt(pt_to_emu(18.0)) == pytest.approx(18.0)
        assert emu_to_pt(914400) == pytest.approx(72.0)

    def test_ideograph_detection_excludes_punctuation(self):
        assert is_ideograph("研")
        assert not is_ideograph("，")
        assert not is_ideograph("A")


# --------------------------------------------------------------------------- #
# Line breaking
# --------------------------------------------------------------------------- #


class TestWrapping:
    def test_cjk_breaks_between_characters(self, font):
        lines = wrap_text("研究方法与实验设计分析", width_pt=60, font_size_pt=20, font=font)
        assert len(lines) > 1
        assert "".join(lines) == "研究方法与实验设计分析"

    def test_latin_breaks_between_words(self, font):
        lines = wrap_text(
            "the proposed method achieves consistent gains",
            width_pt=120,
            font_size_pt=18,
            font=font,
        )
        assert len(lines) > 1
        # No word is split across lines.
        for line in lines:
            assert not line.startswith(" ")
        assert "".join(lines).replace(" ", "") == (
            "theproposedmethodachievesconsistentgains"
        )

    def test_closing_punctuation_never_starts_a_line(self, font):
        """Kinsoku shori: a line beginning with a comma is immediately visible."""
        for width in (40, 60, 80, 100, 140):
            lines = wrap_text(
                "这是一个测试，用于验证标点的换行规则。再加一句，确认效果。",
                width_pt=width,
                font_size_pt=20,
                font=font,
            )
            for line in lines:
                assert line[0] not in "，。、；：？！）」"

    def test_opening_punctuation_never_ends_a_line(self, font):
        for width in (40, 60, 80, 120):
            lines = wrap_text(
                "研究表明（在低资源设定下）模型显著更优（见表一）",
                width_pt=width,
                font_size_pt=20,
                font=font,
            )
            for line in lines:
                assert line[-1] not in "（「《"

    def test_token_wider_than_the_box_is_broken(self, font):
        lines = wrap_text(
            "Pseudoantidisestablishmentarianism",
            width_pt=40,
            font_size_pt=18,
            font=font,
        )
        assert len(lines) > 1

    def test_empty_text_yields_one_empty_line(self, font):
        assert wrap_text("", width_pt=100, font_size_pt=18, font=font) == [""]


# --------------------------------------------------------------------------- #
# Capacity: the headline invariant
# --------------------------------------------------------------------------- #


class TestCapacity:
    GEOMETRIES: ClassVar[list[tuple[float, float, float]]] = [
        (9.0, 3.5, 20.0),
        (9.0, 3.5, 28.0),
        (6.0, 4.0, 14.0),
        (4.0, 2.0, 18.0),
        (3.0, 1.5, 16.0),
        (8.0, 5.0, 12.0),
        (2.0, 1.0, 14.0),
        (5.0, 3.0, 24.0),
    ]

    @pytest.mark.parametrize("w_in,h_in,size", GEOMETRIES)
    @pytest.mark.parametrize("language", ["zh", "en"])
    def test_capacity_is_a_sound_upper_bound(self, font, w_in, h_in, size, language):
        """Text filled to capacity must FIT.

        The planner treats this number as a hard constraint, so an optimistic
        capacity silently produces overflowing slides -- the exact failure this
        whole layer exists to prevent.
        """
        w, h = w_in * 72, h_in * 72
        cap = capacity_units(
            width_pt=w,
            height_pt=h,
            font_size_pt=size,
            font=font,
            language=language,
        )
        if cap == 0:
            return
        m = measure_text(
            _filler_of(cap, language),
            width_pt=w,
            height_pt=h,
            font_size_pt=size,
            font=font,
        )
        assert m.fits, f"capacity {cap} overflows at ratio {m.overflow_ratio:.3f}"

    def test_capacity_shrinks_as_font_grows(self, font):
        caps = [
            capacity_units(
                width_pt=648, height_pt=252, font_size_pt=s, font=font, language="zh"
            )
            for s in (14, 18, 24, 32, 44)
        ]
        assert caps == sorted(caps, reverse=True)

    def test_capacity_grows_with_area(self, font):
        small = capacity_units(
            width_pt=200, height_pt=100, font_size_pt=18, font=font, language="zh"
        )
        large = capacity_units(
            width_pt=600, height_pt=300, font_size_pt=18, font=font, language="zh"
        )
        assert large > small * 3

    def test_degenerate_geometry_is_zero(self, font):
        assert (
            capacity_units(width_pt=0, height_pt=100, font_size_pt=18, font=font) == 0
        )
        assert (
            capacity_units(width_pt=100, height_pt=2, font_size_pt=44, font=font) == 0
        )

    def test_bullet_indent_reduces_capacity(self, font):
        plain = capacity_units(
            width_pt=400, height_pt=200, font_size_pt=18, font=font, language="zh"
        )
        indented = capacity_units(
            width_pt=400,
            height_pt=200,
            font_size_pt=18,
            font=font,
            language="zh",
            bullet_indent_pt=36.0,
        )
        assert indented < plain

    def test_english_capacity_is_in_words(self, font):
        zh = capacity_units(
            width_pt=400, height_pt=200, font_size_pt=18, font=font, language="zh"
        )
        en = capacity_units(
            width_pt=400, height_pt=200, font_size_pt=18, font=font, language="en"
        )
        assert en < zh  # words are coarser than ideographs

    def test_lines_available_is_monotonic(self):
        assert lines_available(height_pt=252, font_size_pt=18) > lines_available(
            height_pt=252, font_size_pt=32
        )
        assert lines_available(height_pt=0, font_size_pt=18) == 0

    def test_overflow_ratio_is_reported_continuously(self, font):
        """The CI metric needs a magnitude, not a boolean."""
        m = measure_text(
            "研" * 400, width_pt=300, height_pt=100, font_size_pt=20, font=font
        )
        assert not m.fits
        assert m.overflow_ratio > 1.5


# --------------------------------------------------------------------------- #
# Template ingestion
# --------------------------------------------------------------------------- #


class TestIngestion:
    def test_stock_template_yields_usable_layouts(self, template):
        contract, report = template
        assert report.layouts_kept >= 8
        assert contract.layouts

    def test_geometry_inheritance_is_resolved(self, template):
        """Layout placeholders usually inherit their box from the master.

        Reading ``ph.left`` naively returns None and silently produces
        zero-sized slots, which would make every capacity zero.
        """
        _, report = template
        assert report.unmeasured_slots == 0

    def test_every_text_slot_has_a_positive_capacity(self, template):
        contract, _ = template
        for layout in contract.layouts:
            for slot in layout.slots:
                if slot.kind in {"title", "subtitle", "body"}:
                    assert slot.capacity_units, (
                        f"{layout.layout_id}/{slot.slot_id} measured zero; a short "
                        "title bar must not be measured at the default 44pt"
                    )

    def test_short_title_bar_gets_a_derived_font_size(self, template):
        """A deliberately short title box must not measure as holding nothing."""
        contract, _ = template
        layout = next(
            (x for x in contract.layouts if "picture" in x.layout_id), None
        )
        if layout is None:
            pytest.skip("stock template has no picture layout")
        title = next(s for s in layout.slots if s.kind == "title")
        assert title.font_size_pt < 44.0
        assert title.capacity_units

    def test_roles_cover_the_academic_essentials(self, template):
        _, report = template
        for role in (
            SlideRole.COVER,
            SlideRole.SECTION,
            SlideRole.TEXT_DENSE,
            SlideRole.TEXT_FIGURE,
            SlideRole.BIG_FIGURE,
        ):
            assert role in report.roles_covered

    def test_chrome_placeholders_are_excluded(self, template):
        """Date, footer and page-number belong to the template, not content."""
        contract, _ = template
        for layout in contract.layouts:
            assert all(s.kind != "footer" for s in layout.slots)

    def test_legibility_floor_is_recorded(self, template):
        contract, _ = template
        body = [
            s
            for layout in contract.layouts
            for s in layout.slots
            if s.kind == "body"
        ]
        assert body
        assert all(s.min_font_size_pt and s.min_font_size_pt >= 14.0 for s in body)

    def test_aspect_ratio_detected(self, template):
        contract, _ = template
        assert contract.page_aspect in {"4:3", "16:9", "16:10", "3:2"}

    def test_missing_font_is_warned_not_fatal(self, tmp_path):
        from pptx import Presentation

        p = tmp_path / "t.pptx"
        Presentation().save(str(p))
        contract, report = ingest_pptx(str(p))
        assert contract.layouts
        assert any("no font" in w for w in report.warnings)


# --------------------------------------------------------------------------- #
# Binding
# --------------------------------------------------------------------------- #


class TestBinding:
    def test_every_slide_gets_a_layout(self, template):
        contract, _ = template
        deck = make_deck(
            [
                Slide(role=SlideRole.COVER, title="封面"),
                Slide(role=SlideRole.TEXT_DENSE, title="要点", blocks=[bullet("一")]),
                Slide(role=SlideRole.SUMMARY, title="总结", blocks=[bullet("二")]),
            ]
        )
        report = bind(deck, contract)
        assert len(report.bindings) == 3
        assert report.ok

    def test_unsupported_role_falls_back_and_warns(self, template):
        """Degrade gracefully; never silently mangle."""
        contract, _ = template
        deck = make_deck(
            [Slide(role=SlideRole.TIMELINE, title="计划", blocks=[bullet("x")])]
        )
        report = bind(deck, contract)
        binding = report.for_slide(deck.sections[0].slides[0].uid)
        assert binding is not None
        if binding.substituted:
            assert report.warnings

    def test_layout_hint_overrides_role_matching(self, template):
        contract, _ = template
        target = contract.layouts[-1].layout_id
        deck = make_deck(
            [
                Slide(
                    role=SlideRole.TEXT_DENSE,
                    title="T",
                    layout_hint=target,
                    blocks=[bullet("x")],
                )
            ]
        )
        report = bind(deck, contract)
        assert report.for_slide(deck.sections[0].slides[0].uid).layout_id == target

    def test_figure_slide_prefers_a_layout_holding_text_too(self, template):
        """A figure without its takeaway line is half a slide."""
        contract, _ = template
        asset = Asset(kind=AssetKind.FIGURE, path="/tmp/x.png")
        slide = Slide(
            role=SlideRole.TEXT_FIGURE,
            title="总体思路",
            blocks=[
                bullet("以结构先验约束表征空间"),
                Block(
                    type=BlockType.FIGURE,
                    asset_ref=asset.uid,
                    spans=[SourceSpan(doc_id=DOC, start=0, end=10)],
                ),
            ],
        )
        deck = make_deck([slide])
        deck.assets.append(asset)
        report = bind(deck, contract)
        layout_id = report.for_slide(slide.uid).layout_id
        layout = next(x for x in contract.layouts if x.layout_id == layout_id)
        regions = sum(
            1 for s in layout.slots if s.kind in {"body", "figure", "table"}
        )
        assert regions >= 2

    def test_vertical_layouts_are_deprioritised(self, template):
        """Chinese academic decks are set horizontally."""
        contract, _ = template
        deck = make_deck(
            [Slide(role=SlideRole.TEXT_DENSE, title="要点", blocks=[bullet("一")])]
        )
        report = bind(deck, contract)
        chosen = report.for_slide(deck.sections[0].slides[0].uid).layout_id
        assert "vertical" not in chosen

    def test_capacity_caps_feed_the_allocator(self, template):
        contract, _ = template
        deck = make_deck(
            [Slide(role=SlideRole.TEXT_DENSE, title="T", blocks=[bullet("一")])]
        )
        report = bind(deck, contract)
        caps = capacity_caps(report)
        assert caps
        for text_cap, lines in caps.values():
            assert text_cap > 0
            assert lines > 0

    def test_geometry_outranks_the_clock(self, template):
        """A generous time budget must not exceed measured capacity."""
        contract, _ = template
        slide = Slide(
            role=SlideRole.TEXT_DENSE,
            title="T",
            blocks=[bullet(f"点{i}", 0.8) for i in range(6)],
        )
        # A very long slot with a tiny deck: time budget would be huge.
        deck = make_deck([slide], seconds=3600, density=Density.COMPACT)
        report = bind(deck, contract)
        caps = capacity_caps(report)
        plan = allocate(deck, capacity=caps)
        measured = caps[slide.uid][0]
        assert plan.slides[0].slide_units <= measured

    def test_capacity_spill_becomes_narration(self, template):
        """Clamping the page lengthens the script; content is never lost."""
        contract, _ = template
        slide = Slide(
            role=SlideRole.TEXT_DENSE,
            title="T",
            blocks=[bullet(f"点{i}", 0.9) for i in range(4)],
        )
        deck = make_deck([slide], seconds=3600, density=Density.COMPACT)
        report = bind(deck, contract)
        plan = allocate(deck, capacity=capacity_caps(report))
        assert plan.slides[0].script_units > 0


# --------------------------------------------------------------------------- #
# Rendering
# --------------------------------------------------------------------------- #


class TestRendering:
    def _deck_with_everything(self, figure_path: str) -> Deck:
        fig = Asset(
            kind=AssetKind.FIGURE,
            path=figure_path,
            intrinsic_aspect=1.6,
            spans=[SourceSpan(doc_id=DOC, start=0, end=10)],
        )
        tbl = Asset(
            kind=AssetKind.TABLE,
            data={
                "columns": ["Method", "Acc"],
                "rows": [["Baseline", 0.812], ["Ours", 0.887]],
            },
            spans=[SourceSpan(doc_id=DOC, start=10, end=20)],
        )
        deck = make_deck(
            [
                Slide(role=SlideRole.COVER, title="标题", subtitle="副标题"),
                Slide(
                    role=SlideRole.TEXT_DENSE,
                    title="三点创新",
                    blocks=[
                        bullet("提出结构一致性正则", 0.95),
                        bullet("参数量降低54%", 0.9),
                        bullet("统一框架", 0.8),
                        Block(
                            type=BlockType.NOTE,
                            content="举医学影像的例子",
                            channel=Channel.SCRIPT,
                        ),
                    ],
                ),
                Slide(
                    role=SlideRole.TEXT_FIGURE,
                    title="总体思路",
                    blocks=[
                        bullet("以结构先验约束表征空间"),
                        Block(
                            type=BlockType.FIGURE,
                            asset_ref=fig.uid,
                            spans=[SourceSpan(doc_id=DOC, start=0, end=10)],
                        ),
                    ],
                ),
                Slide(
                    role=SlideRole.TABLE,
                    title="对比结果",
                    blocks=[
                        Block(
                            type=BlockType.TABLE,
                            asset_ref=tbl.uid,
                            spans=[SourceSpan(doc_id=DOC, start=10, end=20)],
                        )
                    ],
                ),
                Slide(
                    role=SlideRole.EQUATION,
                    title="备用推导",
                    is_backup=True,
                    blocks=[
                        Block(
                            type=BlockType.EQUATION,
                            content=r"\mathcal{R}(h)\le\hat{\mathcal{R}}(h)",
                            spans=[SourceSpan(doc_id=DOC, start=0, end=10)],
                        )
                    ],
                ),
            ]
        )
        deck.assets.extend([fig, tbl])
        return deck

    def _render(self, deck, template, tmp_path, **kw):
        from rostrum.render import render_pptx

        contract, _ = template
        binding = bind(deck, contract)
        allocate(deck, capacity=capacity_caps(binding))
        out = tmp_path / "out.pptx"
        report = render_pptx(
            deck, contract, binding, str(out), font_path=CJK_FONT, **kw
        )
        return report, str(out), binding

    def test_renders_without_overflow(self, template, tmp_path, figure):
        """The headline CI metric."""
        deck = self._deck_with_everything(figure)
        report, _, binding = self._render(deck, template, tmp_path)
        assert report.slides_written == 5
        assert report.overflow_rate == 0.0
        assert not report.missing_assets
        assert overflow_rate(deck, binding) == 0.0

    def test_output_is_natively_editable(self, template, tmp_path, figure):
        """Real text frames and a real table, not a picture of a slide."""
        from pptx import Presentation

        deck = self._deck_with_everything(figure)
        _, out, _ = self._render(deck, template, tmp_path)
        prs = Presentation(out)
        assert len(prs.slides) == 5
        texts = [
            sh.text_frame.text
            for s in prs.slides
            for sh in s.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()
        ]
        assert any("三点创新" in t for t in texts)
        assert any("结构一致性正则" in t for t in texts)
        assert any(sh.has_table for s in prs.slides for sh in s.shapes)
        assert any(sh.shape_type == 13 for s in prs.slides for sh in s.shapes)

    def test_template_layouts_are_reused_not_rebuilt(self, template, tmp_path, figure):
        """Editing preserves theme, fonts and master graphics."""
        from pptx import Presentation

        contract, _ = template
        deck = self._deck_with_everything(figure)
        _, out, _ = self._render(deck, template, tmp_path)
        prs = Presentation(out)
        native_names = {x.native_ref for x in contract.layouts}
        for slide in prs.slides:
            assert slide.slide_layout.name in native_names

    def test_backup_slides_are_moved_to_the_end(self, template, tmp_path, figure):
        """A reserve slide mid-deck gets projected by accident."""
        from pptx import Presentation

        deck = self._deck_with_everything(figure)
        _, out, _ = self._render(deck, template, tmp_path)
        prs = Presentation(out)
        titles = [
            (s.shapes.title.text if s.shapes.title else "") for s in prs.slides
        ]
        assert titles[-1] == "备用推导"

    def test_backup_can_be_excluded(self, template, tmp_path, figure):
        deck = self._deck_with_everything(figure)
        report, _, _ = self._render(
            deck, template, tmp_path, include_backup=False
        )
        assert report.slides_written == 4

    def test_parallel_points_stay_in_one_column(self, template, tmp_path, figure):
        """Splitting three claims across two boxes breaks the reading order."""
        from pptx import Presentation

        deck = self._deck_with_everything(figure)
        _, out, _ = self._render(deck, template, tmp_path)
        prs = Presentation(out)
        target = next(
            s
            for s in prs.slides
            if s.shapes.title and s.shapes.title.text == "三点创新"
        )
        bodies = [
            sh
            for sh in target.shapes
            if sh.has_text_frame
            and sh != target.shapes.title
            and sh.text_frame.text.strip()
        ]
        assert len(bodies) == 1, "parallel bullets were split across columns"
        assert len(bodies[0].text_frame.paragraphs) == 3

    def test_figure_and_its_text_coexist(self, template, tmp_path, figure):
        from pptx import Presentation

        deck = self._deck_with_everything(figure)
        _, out, _ = self._render(deck, template, tmp_path)
        prs = Presentation(out)
        target = next(
            s
            for s in prs.slides
            if s.shapes.title and s.shapes.title.text == "总体思路"
        )
        has_pic = any(sh.shape_type == 13 for sh in target.shapes)
        has_text = any(
            sh.has_text_frame
            and sh != target.shapes.title
            and sh.text_frame.text.strip()
            for sh in target.shapes
        )
        assert has_pic and has_text

    def test_figure_aspect_ratio_is_preserved(self, template, tmp_path, figure):
        """Stretching an author's figure is a visible defect."""
        from pptx import Presentation

        deck = self._deck_with_everything(figure)
        _, out, _ = self._render(deck, template, tmp_path)
        prs = Presentation(out)
        pics = [
            sh for s in prs.slides for sh in s.shapes if sh.shape_type == 13
        ]
        assert pics
        assert pics[0].width / pics[0].height == pytest.approx(1.6, rel=0.02)

    def test_script_channel_lands_in_speaker_notes(self, template, tmp_path, figure):
        """The script is a projection of the same tree, not a second document."""
        from pptx import Presentation

        deck = self._deck_with_everything(figure)
        _, out, _ = self._render(deck, template, tmp_path)
        prs = Presentation(out)
        notes = " ".join(
            s.notes_slide.notes_text_frame.text
            for s in prs.slides
            if s.has_notes_slide
        )
        assert "医学影像" in notes

    def test_missing_asset_is_reported_not_fatal(self, template, tmp_path):
        deck = self._deck_with_everything("/nonexistent/fig.png")
        report, _, _ = self._render(deck, template, tmp_path)
        assert report.missing_assets
        assert report.slides_written == 5  # the deck still renders

    def test_shrinking_stops_at_the_legibility_floor(self, template, tmp_path):
        """An unreadable slide is not a fixed slide."""
        deck = make_deck(
            [
                Slide(
                    role=SlideRole.TEXT_DENSE,
                    title="标题" * 40,
                    blocks=[bullet("要点" * 400)],
                )
            ]
        )
        report, _, _ = self._render(deck, template, tmp_path)
        # It cannot fit, so it must be reported rather than shrunk to 6pt.
        assert report.overflowed_slots
        for _, _, ratio in report.overflowed_slots:
            assert ratio > 1.0

    def test_script_export_is_a_projection(self, template, tmp_path, figure):
        from rostrum.render import export_script

        deck = self._deck_with_everything(figure)
        self._render(deck, template, tmp_path)
        text = export_script(deck, str(tmp_path / "script.md"))
        assert "医学影像" in text
        assert deck.meta.title in text
        assert "备用页" in text  # reserve slides listed separately
