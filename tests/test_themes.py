"""Tests for the built-in themes.

Two things are checked here that ordinary unit tests tend to miss.

**Legibility is measured, not asserted by eye.** Contrast ratios are computed and
compared against the WCAG AA floor, because a theme that looks fine on a laptop
can be unreadable projected in a lit room. A gold accent that looked correct
measured 3.0:1 and was caught only by this check.

**Generated templates are parsed back.** The generator writes OOXML by hand, so
the only real proof it produced a valid package is that python-pptx, the
measurement layer and the binder all consume it and agree on what they found.
"""

from __future__ import annotations

import pytest

from rostrum.ir.enums import Density, Scenario, SlideRole
from rostrum.themes import (
    DEFAULT_THEME_ID,
    REQUIRED_ROLES,
    build_template,
    contrast_ratio,
    get_theme,
    list_themes,
)
from rostrum.themes.spec import Geometry, TypeScale

pytest.importorskip("pptx")


THEME_IDS = [t.theme_id for t in list_themes()]


# --------------------------------------------------------------------------- #
# Specifications
# --------------------------------------------------------------------------- #


def test_there_is_a_default_theme_and_it_exists():
    assert DEFAULT_THEME_ID in THEME_IDS


@pytest.mark.parametrize("theme_id", THEME_IDS)
def test_body_text_clears_the_aa_contrast_floor(theme_id):
    """Body text must be readable from the back of a room, not merely visible."""
    p = get_theme(theme_id).palette
    assert contrast_ratio(p.body, p.background) >= 4.5


@pytest.mark.parametrize("theme_id", THEME_IDS)
def test_accent_clears_the_aa_contrast_floor(theme_id):
    """The accent carries section numbers and emphasis, so it carries words.

    Checked separately from ``body``: an accent chosen for how it looks against a
    background can easily fail as text, which is what happened to this project's
    first gold.
    """
    p = get_theme(theme_id).palette
    assert contrast_ratio(p.accent, p.background) >= 4.5


@pytest.mark.parametrize("theme_id", THEME_IDS)
def test_titles_are_larger_than_body(theme_id):
    ts = get_theme(theme_id).type_scale
    assert ts.title > ts.body > ts.caption


@pytest.mark.parametrize("theme_id", THEME_IDS)
def test_body_size_respects_the_projection_floor(theme_id):
    """Published guidance puts defence body text at 18pt and above.

    ``body_min`` is what the renderer may shrink to; below about 16pt a slide
    stops being readable at the back of a lecture hall, so content must leave the
    page instead.
    """
    ts = get_theme(theme_id).type_scale
    assert ts.body >= 18
    assert ts.body_min >= 16


@pytest.mark.parametrize("theme_id", THEME_IDS)
def test_geometry_leaves_real_margins(theme_id):
    g = get_theme(theme_id).geometry
    # A slide with no margin reads as a word-processor page.
    assert g.margin_left >= 0.05
    assert g.margin_bottom >= 0.05
    # The left gutter is wider than the right: symmetry looks accidental.
    assert g.margin_left >= g.margin_right
    assert g.body_height() > 0.5


def test_type_scale_rejects_an_inverted_hierarchy():
    """A scale whose title is no larger than its body has no hierarchy."""
    with pytest.raises(ValueError, match="larger than body"):
        TypeScale(
            deck_title=40, section=32, title=18, body=20, body_min=16, sub=18,
            caption=14, title_font="f", body_font="f",
            latin_title_font="f", latin_body_font="f",
        )


def test_type_scale_rejects_a_floor_above_the_body_size():
    with pytest.raises(ValueError, match="body_min"):
        TypeScale(
            deck_title=40, section=32, title=28, body=18, body_min=24, sub=18,
            caption=14, title_font="f", body_font="f",
            latin_title_font="f", latin_body_font="f",
        )


def test_geometry_rejects_a_layout_with_no_room_for_content():
    with pytest.raises(ValueError, match="too little room"):
        Geometry(
            margin_left=0.1, margin_right=0.1, margin_top=0.3,
            margin_bottom=0.35, title_height=0.3, gap=0.05,
        )


def test_unknown_theme_names_the_available_ones():
    with pytest.raises(ValueError, match="available"):
        get_theme("no-such-theme")


# --------------------------------------------------------------------------- #
# Generated packages
# --------------------------------------------------------------------------- #


@pytest.fixture(scope="module", params=THEME_IDS)
def built(request, tmp_path_factory):
    theme = get_theme(request.param)
    out = tmp_path_factory.mktemp("themes") / f"{theme.theme_id}.pptx"
    build_template(theme, str(out))
    return theme, str(out)


def test_generated_template_opens_as_a_presentation(built):
    """The generator writes OOXML by hand; this is the proof it is valid."""
    import pptx

    theme, path = built
    prs = pptx.Presentation(path)
    assert len(prs.slide_masters) == 1
    assert len(prs.slide_layouts) >= 12
    assert (prs.slide_width, prs.slide_height) == theme.slide_size_emu()


def test_generated_template_covers_every_required_role(built):
    """A missing role forces the binder to substitute, which looks improvised."""
    from rostrum.templates import ingest_pptx

    theme, path = built
    contract, _ = ingest_pptx(path, template_id=theme.theme_id, license="builtin")
    covered = contract.supported_roles()
    missing = [r.value for r in REQUIRED_ROLES if r not in covered]
    assert missing == []


def test_measured_sizes_match_the_declared_type_scale(built):
    """Declared sizes must be readable by the measurement layer.

    The first build wrote sizes only into a shape-level ``lstStyle``, which
    python-pptx cannot see; every title then measured at PowerPoint's 44pt
    default and capacities were computed against the wrong size.
    """
    from rostrum.templates import ingest_pptx

    theme, path = built
    contract, _ = ingest_pptx(path, template_id=theme.theme_id, license="builtin")

    content = next(x for x in contract.layouts if x.layout_id == "title-and-content")
    title = next(s for s in content.slots if s.kind == "title")
    body = next(s for s in content.slots if s.kind == "body")

    assert title.font_size_pt == pytest.approx(theme.type_scale.title, abs=0.6)
    assert body.font_size_pt == pytest.approx(theme.type_scale.body, abs=0.6)


def test_theme_fonts_are_discoverable_including_cjk(built):
    """CJK faces live in the theme part, not on shapes.

    A template that references ``+mj-lt`` and declares real faces once in the
    theme is correctly built -- and reported no fonts at all until the ingest
    layer learned to read the theme part.
    """
    from rostrum.templates import ingest_pptx

    theme, path = built
    contract, _ = ingest_pptx(path, template_id=theme.theme_id, license="builtin")
    assert contract.fonts
    assert any("CJK" in f or "YaHei" in f or "Hei" in f for f in contract.fonts)


def test_no_layout_titles_the_page_below_the_fold(built):
    """A built-in theme must never need the low-title veto to save it."""
    from rostrum.templates import ingest_pptx

    theme, path = built
    contract, _ = ingest_pptx(path, template_id=theme.theme_id, license="builtin")
    for layout in contract.layouts:
        titles = [s for s in layout.slots if s.kind == "title"]
        if titles:
            assert min(t.box.y for t in titles) <= 0.45, layout.layout_id


# --------------------------------------------------------------------------- #
# End to end through every theme
# --------------------------------------------------------------------------- #


@pytest.fixture(scope="module")
def manuscript(tmp_path_factory):
    import docx
    from docx.shared import Inches
    from PIL import Image

    d = tmp_path_factory.mktemp("src")
    img = d / "f.png"
    Image.new("RGB", (800, 500), "white").save(img)

    doc = docx.Document()
    doc.core_properties.title = "主题渲染测试"
    doc.core_properties.author = "张三"
    doc.add_heading("主题渲染测试", 1)
    doc.add_paragraph("张三　某大学计算机学院")
    doc.add_heading("研究背景与问题", 1)
    doc.add_paragraph("现有方法在小样本条件下性能显著下降，这是本项目要解决的核心问题。")
    doc.add_heading("创新点", 1)
    doc.add_paragraph("精度提升7.5个百分点", style="List Bullet")
    doc.add_paragraph("参数量降低54%", style="List Bullet")
    doc.add_heading("研究目标与内容", 1)
    doc.add_paragraph("总体框架如下图所示。")
    doc.add_picture(str(img), width=Inches(4.0))
    doc.add_paragraph("图1 整体技术框架")
    doc.add_heading("研究基础", 1)
    doc.add_paragraph("表1 对比结果")
    t = doc.add_table(rows=2, cols=3)
    for j, v in enumerate(["方法", "准确率", "参数量"]):
        t.cell(0, j).text = v
    for j, v in enumerate(["本文", "0.887", "11M"]):
        t.cell(1, j).text = v
    doc.add_heading("可行性分析", 1)
    doc.add_paragraph("团队具备完整的算力条件与数据积累。")

    path = d / "m.docx"
    doc.save(str(path))
    return str(path)


def test_every_theme_renders_without_overflow(built, manuscript, tmp_path):
    """The real acceptance test: a manuscript through each theme, nothing spilling."""
    from rostrum.budget import allocate
    from rostrum.ingest.docx_parser import parse_docx
    from rostrum.ingest.planner import plan_deck
    from rostrum.render import render_pptx
    from rostrum.templates import bind, capacity_caps, ingest_pptx, overflow_rate

    theme, path = built
    doc = parse_docx(manuscript, asset_dir=str(tmp_path / "a"))
    deck = plan_deck(doc, total_seconds=480, scenario=Scenario.GRANT_DEFENSE)

    contract, _ = ingest_pptx(path, template_id=theme.theme_id, license="builtin")
    binding = bind(deck, contract)
    allocate(deck, apply=True, capacity=capacity_caps(binding))

    out = tmp_path / f"{theme.theme_id}.pptx"
    report = render_pptx(deck, contract, binding, str(out))

    assert out.exists()
    assert report.slides_written >= 8
    assert report.missing_assets == []
    assert overflow_rate(deck, binding) == 0.0
    # No slide should have needed a substituted layout: a built-in theme covers
    # everything the planner emits.
    assert not any(b.substituted for b in binding.bindings.values())


def test_planner_produces_a_full_talk_structure(manuscript, tmp_path):
    """A talk needs a cover, an agenda and dividers, not just content pages."""
    from rostrum.ingest.docx_parser import parse_docx
    from rostrum.ingest.planner import plan_deck

    doc = parse_docx(manuscript, asset_dir=str(tmp_path / "a"))
    deck = plan_deck(doc, total_seconds=480, scenario=Scenario.GRANT_DEFENSE)
    roles = [slide.role for _, slide in deck.iter_slides()]

    assert roles[0] is SlideRole.COVER
    assert SlideRole.AGENDA in roles
    assert SlideRole.SECTION in roles
    assert roles[-1] is SlideRole.ACKNOWLEDGEMENT


def test_navigation_slides_are_excluded_from_the_content_count(manuscript, tmp_path):
    """Pacing advice is meaningless if dividers count as content.

    A cover plus four dividers is five pages and about ten seconds; counting them
    made a correctly sized deck report as 60% over its slide target.
    """
    from rostrum.ingest.docx_parser import parse_docx
    from rostrum.ingest.planner import plan_deck

    doc = parse_docx(manuscript, asset_dir=str(tmp_path / "a"))
    deck = plan_deck(doc, total_seconds=480, scenario=Scenario.GRANT_DEFENSE)

    total = sum(len(s.slides) for s in deck.sections)
    assert deck.navigation_slide_count > 0
    assert deck.content_slide_count < total
    assert deck.content_slide_count + deck.navigation_slide_count == total


def test_navigation_slides_barely_consume_time(manuscript, tmp_path):
    """A divider is spoken over in a couple of seconds, not ten."""
    from rostrum.budget import allocate
    from rostrum.ingest.docx_parser import parse_docx
    from rostrum.ingest.planner import plan_deck

    doc = parse_docx(manuscript, asset_dir=str(tmp_path / "a"))
    deck = plan_deck(doc, total_seconds=480, scenario=Scenario.GRANT_DEFENSE)
    allocate(deck, apply=True)

    dividers = [
        slide
        for _, slide in deck.iter_slides()
        if slide.role is SlideRole.SECTION
    ]
    assert dividers
    for slide in dividers:
        assert slide.dwell_seconds is not None
        assert slide.dwell_seconds < 8, "a divider must not eat a content slide's time"


def test_agenda_items_cite_the_headings_they_name(manuscript, tmp_path):
    """An agenda is not authored from nothing: each entry is a real heading."""
    from rostrum.ingest.docx_parser import parse_docx
    from rostrum.ingest.planner import plan_deck

    doc = parse_docx(manuscript, asset_dir=str(tmp_path / "a"))
    deck = plan_deck(doc, total_seconds=480, scenario=Scenario.GRANT_DEFENSE)

    agenda = next(
        slide for _, slide in deck.iter_slides() if slide.role is SlideRole.AGENDA
    )
    assert agenda.blocks
    for block in agenda.blocks:
        assert block.spans, "an agenda item must cite the heading it names"
        assert doc.text[block.spans[0].start : block.spans[0].end] == block.content


@pytest.mark.parametrize("density", [Density.SPARSE, Density.COMPACT])
def test_density_still_fits_the_built_in_default(manuscript, tmp_path, density):
    from rostrum.budget import allocate
    from rostrum.ingest.docx_parser import parse_docx
    from rostrum.ingest.planner import plan_deck
    from rostrum.render import render_pptx
    from rostrum.templates import bind, capacity_caps, ingest_pptx, overflow_rate

    template = tmp_path / "t.pptx"
    build_template(get_theme(DEFAULT_THEME_ID), str(template))

    doc = parse_docx(manuscript, asset_dir=str(tmp_path / "a"))
    deck = plan_deck(doc, total_seconds=480, density=density)
    contract, _ = ingest_pptx(str(template), template_id="t", license="builtin")
    binding = bind(deck, contract)
    allocate(deck, apply=True, capacity=capacity_caps(binding))

    out = tmp_path / f"{density.value}.pptx"
    render_pptx(deck, contract, binding, str(out))
    assert overflow_rate(deck, binding) == 0.0
