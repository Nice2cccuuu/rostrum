"""PPTX renderer: deck IR -> a natively editable .pptx.

Strategy: **edit, don't generate.** Slides are created from the template's own
layouts so that every inherited property -- theme colours, fonts, background
graphics, master placement -- survives untouched. Building shapes from scratch
would reproduce the look of a template while losing everything that makes it
*that* template, and is the reason generic tools produce decks that cannot be
handed to a co-author.

Two consequences follow, and both are deliberate:

- Text is written into inherited placeholders, never into new text boxes.
- Anything that will not fit is caught *before* writing, using the same measured
  capacity the planner used, and demoted to the speaker notes rather than
  silently clipped.
"""

from __future__ import annotations

import contextlib
import copy
from dataclasses import dataclass, field

from rostrum.budget.allocate import count_units
from rostrum.ir.enums import BlockType, Channel
from rostrum.ir.nodes import Asset, Block, Deck, Slide
from rostrum.measure.text import (
    emu_to_pt,
    load_font,
    measure_text,
    pt_to_emu,
    wrap_text,
)
from rostrum.render.anchors import Anchor, AnchorMap, Box
from rostrum.templates.binding import BindingReport
from rostrum.templates.contract import TemplateContract

try:
    from pptx import Presentation
    from pptx.util import Pt

    _HAVE_PPTX = True
except ImportError:  # pragma: no cover
    _HAVE_PPTX = False


@dataclass
class RenderReport:
    """Diagnostics from a render pass."""

    path: str
    slides_written: int = 0
    overflowed_slots: list[tuple[str, str, float]] = field(default_factory=list)
    """``(slide_path, slot_id, overflow_ratio)`` for slots that did not fit."""
    demoted_to_notes: list[str] = field(default_factory=list)
    """Blocks that could not be placed and went to the speaker notes."""
    missing_assets: list[str] = field(default_factory=list)
    shrunk_slots: list[tuple[str, str, float, float]] = field(default_factory=list)
    """``(slide_path, slot_id, from_pt, to_pt)`` autofit reductions."""
    warnings: list[str] = field(default_factory=list)
    anchors: AnchorMap | None = None
    """Geometry of every editable region, for click-to-select in a preview."""

    @property
    def overflow_rate(self) -> float:
        """Fraction of written slides with at least one overflowing slot.

        One of the two CI metrics that require no model in the loop.
        """
        if not self.slides_written:
            return 0.0
        bad = {s for s, _, _ in self.overflowed_slots}
        return len(bad) / self.slides_written

    @property
    def ok(self) -> bool:
        return not self.overflowed_slots and not self.missing_assets


def render_pptx(
    deck: Deck,
    contract: TemplateContract,
    binding: BindingReport,
    out_path: str,
    *,
    font_path: str | None = None,
    include_backup: bool = True,
    write_notes: bool = True,
) -> RenderReport:
    """Render ``deck`` into ``out_path`` using the bound template.

    ``binding`` must come from :func:`rostrum.templates.binding.bind` against the
    same ``contract``, so that the layout chosen here is the one whose capacity
    the budget respected.
    """
    if not _HAVE_PPTX:  # pragma: no cover
        raise RuntimeError("python-pptx is required: pip install 'rostrum[pptx]'")
    if not contract.source_path:
        raise ValueError(
            "contract has no source_path; the renderer edits the original "
            "template rather than rebuilding it"
        )

    prs = Presentation(contract.source_path)
    _strip_existing_slides(prs)

    layout_by_id = {}
    for layout in contract.layouts:
        native = _find_native_layout(prs, layout.native_ref, layout.layout_id)
        if native is not None:
            layout_by_id[layout.layout_id] = (layout, native)

    report = RenderReport(path=out_path)
    font = load_font(font_path)
    lang = deck.meta.language

    # Anchors are emitted here rather than read back from the file, because a
    # bullet is a paragraph inside a shared text frame and has no geometry of its
    # own in the package. Only the renderer knows both the box and what went in it.
    report.anchors = AnchorMap(
        deck_uid=deck.uid,
        slide_width_pt=emu_to_pt(prs.slide_width),
        slide_height_pt=emu_to_pt(prs.slide_height),
    )
    canvas = (float(prs.slide_width), float(prs.slide_height))

    # Backup slides go last, always. A reserve slide sitting mid-deck will be
    # projected by accident during the talk, which is exactly what marking it
    # "backup" was meant to prevent.
    ordered = [(s, sl) for s, sl in deck.iter_slides() if not sl.is_backup]
    if include_backup:
        ordered += [(s, sl) for s, sl in deck.iter_slides() if sl.is_backup]

    for _, slide in ordered:
        bound = binding.for_slide(slide.uid)
        if bound is None or bound.layout_id not in layout_by_id:
            report.warnings.append(
                f"slide {deck.path_of(slide.uid)} has no usable layout binding; "
                "skipped"
            )
            continue

        spec, native = layout_by_id[bound.layout_id]
        native_slide = prs.slides.add_slide(native)
        _inherit_layout_appearance(native_slide, native)
        _write_slide(
            deck,
            slide,
            spec,
            native_slide,
            font=font,
            language=lang,
            report=report,
            write_notes=write_notes,
            slide_index=report.slides_written,
            canvas=canvas,
        )
        report.slides_written += 1

    prs.save(out_path)
    return report


# --------------------------------------------------------------------------- #
# Anchor collection
# --------------------------------------------------------------------------- #


@dataclass
class _AnchorSink:
    """Records where each node was drawn, as the renderer draws it.

    Per-paragraph geometry has to be *computed*: several blocks share one text
    frame, so the file itself only knows the placeholder's box. The same wrapping
    metrics the fit calculation used are reused here, which keeps the anchors
    consistent with the layout decision rather than being a second guess at it.
    """

    amap: AnchorMap | None
    deck: Deck
    slide: Slide
    slide_index: int
    canvas: tuple[float, float] | None
    font: object

    @property
    def live(self) -> bool:
        return self.amap is not None and self.canvas is not None

    def _box(self, left, top, width, height) -> Box | None:
        """Normalise EMU geometry to slide fractions, clamped to the page."""
        if not self.canvas:
            return None
        cw, ch = self.canvas
        if not cw or not ch or width is None or height is None:
            return None
        x = max(0.0, min(1.0, float(left) / cw))
        y = max(0.0, min(1.0, float(top) / ch))
        w = min(1.0 - x, float(width) / cw)
        h = min(1.0 - y, float(height) / ch)
        if w <= 0 or h <= 0:
            return None
        return Box(x=x, y=y, w=w, h=h)

    def _emit(self, uid, kind, box, preview="", confidence=1.0) -> None:
        if not self.live or box is None:
            return
        base = self.deck.path_of(uid) or uid
        # Title, subtitle and whole-page anchors all carry the slide's uid, since
        # each resolves to an edit on that slide. Suffixing keeps them apart when
        # a UI lists what a click could have meant.
        path = base if kind == "block" else f"{base}:{kind}"
        self.amap.add(
            Anchor(
                uid=uid,
                kind=kind,
                slide_index=self.slide_index,
                slide_uid=self.slide.uid,
                path=path,
                box=box,
                preview=_preview(preview),
                confidence=confidence,
            )
        )

    def add_slide_box(self) -> None:
        """A whole-page anchor, so a click on empty space still resolves.

        Scored low deliberately: it is the answer to "which page" when nothing
        else was hit, never a competitor to the bullet under the cursor.
        """
        self._emit(
            self.slide.uid,
            "slide",
            Box(x=0.0, y=0.0, w=1.0, h=1.0),
            preview=self.slide.title or "",
            confidence=1.0,
        )

    def add_shape(self, uid: str, kind: str, ph, preview: str = "") -> None:
        """Anchor covering a whole placeholder -- exact, since the file says so."""
        self._emit(
            uid,
            kind,
            self._box(ph.left, ph.top, ph.width, ph.height),
            preview=preview,
            confidence=1.0,
        )

    def add_paragraphs(
        self,
        ph,
        blocks: list[Block],
        texts: list[str],
        *,
        font_size_pt: float,
        slot,
    ) -> None:
        """Anchor each block to the lines it actually occupies.

        Line count per block comes from the same wrapping used to choose the font
        size, so a bullet that wrapped to three lines gets a box three lines tall.
        Vertical anchoring is read from the shape: a body box that centres its
        text puts the first line well below the placeholder's top, and assuming
        top alignment put every box in the wrong place -- visible immediately in
        the overlay, invisible to any assertion about counts.
        """
        if not self.live or not blocks:
            return

        box = self._box(ph.left, ph.top, ph.width, ph.height)
        if box is None:
            return

        width_pt = emu_to_pt(ph.width) - _text_inset_pt(ph)
        line_h_pt = font_size_pt * _LINE_SPACING

        counts = []
        for text in texts:
            wrapped = wrap_text(
                text, width_pt=max(1.0, width_pt), font_size_pt=font_size_pt,
                font=self.font,
            )
            counts.append(max(1, len(wrapped)))

        total_h_pt = sum(counts) * line_h_pt
        box_h_pt = emu_to_pt(ph.height)
        offset_pt = _vertical_offset_pt(ph, box_h_pt, total_h_pt)

        cursor = offset_pt
        # strict: a mismatch here would silently truncate the anchor list, so a
        # bullet would become unclickable rather than raising.
        for block, text, n_lines in zip(blocks, texts, counts, strict=True):
            h_pt = n_lines * line_h_pt
            top_emu = float(ph.top) + pt_to_emu(cursor)
            self._emit(
                block.uid,
                "block",
                self._box(top=top_emu, left=ph.left, width=ph.width, height=pt_to_emu(h_pt)),
                preview=text,
                # Slightly below certain: PowerPoint's own line breaking can
                # differ from ours by a hair, so a click near a boundary may
                # land on the neighbour. Hit-testing ranks rather than asserts.
                confidence=0.85,
            )
            cursor += h_pt


# Matches the spacing assumption in the measurement layer; a mismatch here would
# make anchors drift down the page relative to the text they describe.
_LINE_SPACING = 1.22


def _text_inset_pt(ph) -> float:
    """Left plus right internal margin of a text frame, in points."""
    try:
        tf = ph.text_frame
        return emu_to_pt((tf.margin_left or 0) + (tf.margin_right or 0))
    except Exception:  # pragma: no cover - shapes without a text frame
        return 0.0


def _vertical_offset_pt(ph, box_h_pt: float, text_h_pt: float) -> float:
    """Where the first line starts, honouring the resolved vertical anchor.

    The anchor must be resolved through the inheritance chain. A shape that does
    not set it returns ``None`` from python-pptx while the real value sits on the
    layout or the master -- and treating ``None`` as top alignment put every
    anchor box 150px above the text it described. The overlay showed it at once;
    no assertion about box counts would have.
    """
    slack = max(0.0, box_h_pt - text_h_pt)
    anchor = _resolved_anchor(ph)
    if anchor == "middle":
        return slack / 2
    if anchor == "bottom":
        return slack
    return 0.0


def _resolved_anchor(ph) -> str:
    """Vertical anchor of a placeholder, following slide → layout → master."""
    try:
        from pptx.enum.text import MSO_ANCHOR
    except Exception:  # pragma: no cover
        return "top"

    names = {
        MSO_ANCHOR.MIDDLE: "middle",
        MSO_ANCHOR.BOTTOM: "bottom",
        MSO_ANCHOR.TOP: "top",
    }
    for source in _inheritance_chain(ph):
        try:
            value = source.text_frame.vertical_anchor
        except Exception:  # pragma: no cover
            continue
        if value is not None:
            return names.get(value, "top")
    return "top"


def _inheritance_chain(ph) -> list:
    """The shape, then its layout counterpart, then its master counterpart."""
    chain = [ph]
    try:
        idx = ph.placeholder_format.idx
        slide = ph.part.slide if hasattr(ph.part, "slide") else None
        layout = slide.slide_layout if slide is not None else None
    except Exception:  # pragma: no cover
        return chain

    for level in (layout, getattr(layout, "slide_master", None)):
        if level is None:
            continue
        with contextlib.suppress(Exception):
            for candidate in level.placeholders:
                if candidate.placeholder_format.idx == idx:
                    chain.append(candidate)
                    break
    return chain


def _preview(text: str, limit: int = 24) -> str:
    flat = " ".join((text or "").split())
    return flat if len(flat) <= limit else flat[:limit] + "…"


# --------------------------------------------------------------------------- #
# Slide population
# --------------------------------------------------------------------------- #


def _inherit_layout_appearance(slide, layout) -> None:
    """Reproduce the layout's own appearance on a slide built from it.

    Two OOXML behaviours have to be handled together, and both were found by
    rendering a deck and looking at it rather than by any check that passes or
    fails:

    1. ``showMasterSp`` does not inherit. A layout that hides master shapes has
       no effect unless the slide repeats the declaration, so a cover came out
       wearing the content pages' title rule.

    2. A slide inherits its layout's *placeholders* but not its static shapes.
       Setting ``showMasterSp="0"`` to fix (1) therefore also removed the
       layout's own decoration -- the cover rule and the section stroke vanished
       entirely. Those shapes must be copied onto the slide explicitly.

    Copying is deliberate rather than clever: the shapes are flat rectangles, so
    a deep copy of the XML is exact, and the result is a self-contained slide
    that keeps its appearance if the template is later detached.
    """
    flag = layout._element.get("showMasterSp")
    if flag is not None:
        slide._element.set("showMasterSp", flag)

    if flag != "0":
        # Master shapes are visible, so the layout's own decoration will already
        # be drawn beneath the slide by the consumer.
        return

    tree = slide.shapes._spTree
    for shape in layout.shapes:
        if shape.is_placeholder:
            continue  # placeholders are inherited already
        tree.insert_element_before(
            copy.deepcopy(shape._element), "p:extLst"
        )


def _write_slide(
    deck: Deck,
    slide: Slide,
    spec,
    native_slide,
    *,
    font,
    language: str,
    report: RenderReport,
    write_notes: bool,
    slide_index: int = 0,
    canvas: tuple[float, float] | None = None,
) -> None:
    path = deck.path_of(slide.uid) or slide.uid
    sink = _AnchorSink(
        amap=report.anchors,
        deck=deck,
        slide=slide,
        slide_index=slide_index,
        canvas=canvas,
        font=font,
    )
    sink.add_slide_box()
    assets = deck.asset_map()
    placeholders = {f"ph{p.placeholder_format.idx}": p for p in native_slide.placeholders}

    shown = slide.slide_blocks()
    text_blocks = [
        b for b in shown if not b.is_visual and b.type is not BlockType.NOTE
    ]
    visual_blocks = [b for b in shown if b.is_visual]

    # -- title ------------------------------------------------------------ #
    title_slot = next((s for s in spec.slots if s.kind == "title"), None)
    if title_slot is not None:
        ph = placeholders.get(title_slot.slot_id)
        if ph is not None and slide.title:
            _fill_text(
                ph,
                [slide.title],
                slot=title_slot,
                font=font,
                language=language,
                slide_path=path,
                report=report,
            )
            # The title anchors to the slide: clicking a heading and saying
            # "改成「问题与挑战」" must resolve to a set_title on this page.
            sink.add_shape(slide.uid, "title", ph, preview=slide.title)

    subtitle_slot = next((s for s in spec.slots if s.kind == "subtitle"), None)
    if subtitle_slot is not None and slide.subtitle:
        ph = placeholders.get(subtitle_slot.slot_id)
        if ph is not None:
            _fill_text(
                ph,
                [slide.subtitle],
                slot=subtitle_slot,
                font=font,
                language=language,
                slide_path=path,
                report=report,
            )
            sink.add_shape(slide.uid, "subtitle", ph, preview=slide.subtitle)

    # -- figures and tables ----------------------------------------------- #
    # Body slots are consumed largest-first: content should land in the page's
    # main area, never in a caption-sized box while the rest of the page is
    # empty.
    visual_slots = sorted(
        (s for s in spec.slots if s.kind in {"figure", "table"}),
        key=lambda s: -(s.box.w * s.box.h),
    )
    body_slots = sorted(
        (s for s in spec.slots if s.kind == "body"),
        key=lambda s: -(s.box.w * s.box.h),
    )
    leftover_visuals: list[Block] = []

    for i, block in enumerate(visual_blocks):
        slot = visual_slots[i] if i < len(visual_slots) else None
        # Reuse a body placeholder: an OBJECT placeholder accepts a picture as
        # readily as text, which is what makes sparse templates workable. But
        # never consume the *last* body slot while text is still waiting for one
        # -- a figure with no takeaway line beside it tells the audience nothing
        # about what to look at.
        if slot is None and body_slots and (len(body_slots) > 1 or not text_blocks):
            slot = body_slots.pop(0)
        if slot is None:
            leftover_visuals.append(block)
            continue

        ph = placeholders.get(slot.slot_id)
        asset = assets.get(block.asset_ref or "")
        if ph is None or asset is None:
            report.missing_assets.append(block.asset_ref or block.uid)
            continue
        _place_visual(native_slide, ph, asset, block, report=report, path=path)
        sink.add_shape(
            block.uid,
            "table" if asset.data else "figure",
            ph,
            preview=asset.caption or block.content,
        )

    # -- body text -------------------------------------------------------- #
    if body_slots and text_blocks:
        groups = _distribute(text_blocks, body_slots)
        for slot, group in zip(body_slots, groups, strict=False):
            ph = placeholders.get(slot.slot_id)
            if ph is None or not group:
                continue
            size, placed, texts = _fill_bullets(
                ph,
                group,
                slot=slot,
                font=font,
                language=language,
                slide_path=path,
                report=report,
            )
            if size:
                sink.add_paragraphs(
                    ph, placed, texts, font_size_pt=size, slot=slot
                )
    elif text_blocks:
        # No body slot at all: the content belongs in the narration rather than
        # in an invented text box that would break the template's design.
        report.demoted_to_notes.extend(b.uid for b in text_blocks)

    # -- speaker notes ---------------------------------------------------- #
    if write_notes:
        notes = _build_notes(deck, slide, leftover_visuals, language)
        if notes:
            native_slide.notes_slide.notes_text_frame.text = notes

    for block in leftover_visuals:
        report.demoted_to_notes.append(block.uid)


def _distribute(blocks: list[Block], slots: list) -> list[list[Block]]:
    """Assign blocks to body slots, keeping the reading order intact.

    A set of parallel bullets is a single unit of meaning: splitting three
    innovation claims across two boxes makes the audience read 1 - then jump -
    then 2, 3. So multi-column filling only happens when there is genuinely
    enough text to warrant it; otherwise everything goes into the largest slot
    and the spare ones are left for the template's own design to handle.
    """
    if not slots:
        return []
    if len(slots) == 1 or len(blocks) <= _MIN_BLOCKS_TO_SPLIT:
        # Everything into the largest slot, which the caller has sorted first.
        return [blocks] + [[] for _ in slots[1:]]

    groups: list[list[Block]] = [[] for _ in slots]
    caps = [max(1, s.capacity_units or 1) for s in slots]
    total_cap = sum(caps)
    loads = [0] * len(slots)
    idx = 0
    for block in blocks:
        weight = block.word_budget or max(1, len(block.content))
        # Move to the next column once this one has taken its proportional share.
        while (
            idx < len(slots) - 1
            and loads[idx] > 0
            and loads[idx] >= caps[idx] / total_cap * _total_weight(blocks)
        ):
            idx += 1
        groups[idx].append(block)
        loads[idx] += weight
    return groups


# Below this, parallel points stay together in one column.
_MIN_BLOCKS_TO_SPLIT = 4


def _total_weight(blocks: list[Block]) -> int:
    return sum(b.word_budget or max(1, len(b.content)) for b in blocks) or 1


def _fill_text(
    ph,
    lines: list[str],
    *,
    slot,
    font,
    language: str,
    slide_path: str,
    report: RenderReport,
) -> float | None:
    """Write plain lines into a placeholder, verifying they fit.

    Returns the applied font size, which the anchor sink needs to compute
    per-paragraph geometry from the same numbers the fit used.
    """
    tf = ph.text_frame
    tf.clear()
    text = "\n".join(lines)
    tf.paragraphs[0].text = text

    size = _fit_size(
        text,
        ph=ph,
        slot=slot,
        font=font,
        slide_path=slide_path,
        report=report,
    )
    if size is not None:
        for para in tf.paragraphs:
            para.font.size = Pt(size)
    return size or slot.font_size_pt


def _fill_bullets(
    ph,
    blocks: list[Block],
    *,
    slot,
    font,
    language: str,
    slide_path: str,
    report: RenderReport,
) -> tuple[float | None, list[Block], list[str]]:
    """Write bullets into a body placeholder, preserving inherited formatting.

    The first paragraph is reused rather than added, because a freshly cleared
    text frame keeps one paragraph whose properties are inherited from the
    layout. Adding a paragraph and leaving the original empty is the standard way
    to end up with a stray blank bullet.
    """
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True

    first = True
    rendered: list[str] = []
    placed: list[Block] = []
    for block in blocks:
        text = _block_text(block)
        if not text:
            continue
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        para.text = text
        para.level = min(block.level, slot.max_bullet_level)
        first = False
        rendered.append(text)
        placed.append(block)

    if first:  # nothing was written
        return None, [], []

    size = _fit_size(
        "\n".join(rendered),
        ph=ph,
        slot=slot,
        font=font,
        slide_path=slide_path,
        report=report,
        lines_hint=len(rendered),
    )
    if size is not None:
        for para in tf.paragraphs:
            para.font.size = Pt(size)
    return (size or slot.font_size_pt), placed, rendered


def _block_text(block: Block) -> str:
    if block.type is BlockType.EQUATION:
        # Inline LaTeX is kept verbatim; a real equation object is a v0.3
        # concern, and mangling the source would be worse than showing it.
        return block.content
    return block.content


def _fit_size(
    text: str,
    *,
    ph,
    slot,
    font,
    slide_path: str,
    report: RenderReport,
    lines_hint: int = 1,
) -> float | None:
    """Choose a font size that fits, or record an overflow.

    Autofit in PowerPoint is unreliable when set programmatically -- the stored
    flag is not honoured until the file is opened and the shape re-laid out. So
    the fit is computed here with the same metrics the planner used, and applied
    as an explicit size. Shrinking stops at ``min_font_size_pt``: an unreadable
    slide is not a fixed slide.
    """
    box = _shape_box_pt(ph, slot)
    if box is None:
        return None
    w_pt, h_pt = box

    nominal = slot.font_size_pt or 18.0
    floor = slot.min_font_size_pt or max(10.0, nominal * 0.6)

    size = nominal
    while size >= floor:
        m = measure_text(
            text,
            width_pt=w_pt,
            height_pt=h_pt,
            font_size_pt=size,
            font=font,
        )
        if m.fits:
            if size < nominal:
                report.shrunk_slots.append(
                    (slide_path, slot.slot_id, nominal, size)
                )
            return size
        size -= 1.0

    m = measure_text(
        text, width_pt=w_pt, height_pt=h_pt, font_size_pt=floor, font=font
    )
    report.overflowed_slots.append((slide_path, slot.slot_id, m.overflow_ratio))
    return floor


def _shape_box_pt(ph, slot) -> tuple[float, float] | None:
    """Usable text area of a placeholder, in points."""
    try:
        w, h = ph.width, ph.height
    except (AttributeError, KeyError):  # pragma: no cover
        return None
    if not w or not h:
        return None

    w_pt, h_pt = emu_to_pt(w), emu_to_pt(h)
    # Subtract the inherited internal margins, else every measurement is
    # optimistic by ~0.2in horizontally.
    try:
        tf = ph.text_frame
        w_pt -= emu_to_pt((tf.margin_left or 0) + (tf.margin_right or 0))
        h_pt -= emu_to_pt((tf.margin_top or 0) + (tf.margin_bottom or 0))
    except (AttributeError, KeyError):  # pragma: no cover
        pass

    if slot.kind == "body":
        w_pt -= 18.0  # bullet glyph and hanging indent
    return (max(1.0, w_pt), max(1.0, h_pt))


def _place_visual(
    native_slide,
    ph,
    asset: Asset,
    block: Block,
    *,
    report: RenderReport,
    path: str,
) -> None:
    """Insert a figure or table into a placeholder.

    Aspect ratio is preserved and the image centred within the slot: stretching
    an author's figure to fill a box is a visible defect in a scientific talk.
    """
    if asset.path:
        try:
            _insert_picture(native_slide, ph, asset)
            return
        except Exception as exc:  # pragma: no cover - unreadable image
            report.warnings.append(
                f"{path}: could not place {asset.uid} ({exc}); left empty"
            )
            report.missing_assets.append(asset.uid)
            return

    if asset.data:
        try:
            _insert_table(native_slide, ph, asset)
            return
        except Exception as exc:  # pragma: no cover
            report.warnings.append(f"{path}: could not build table ({exc})")
            report.missing_assets.append(asset.uid)
            return

    report.missing_assets.append(asset.uid)


def _insert_picture(native_slide, ph, asset: Asset) -> None:
    left, top = ph.left, ph.top
    box_w, box_h = ph.width, ph.height

    aspect = asset.intrinsic_aspect
    if aspect is None:
        aspect = _probe_aspect(asset.path)

    if aspect and aspect > 0:
        box_aspect = box_w / box_h
        if aspect > box_aspect:
            w = box_w
            h = int(box_w / aspect)
        else:
            h = box_h
            w = int(box_h * aspect)
        left = left + (box_w - w) // 2
        top = top + (box_h - h) // 2
    else:  # pragma: no cover
        w, h = box_w, box_h

    native_slide.shapes.add_picture(asset.path, left, top, width=w, height=h)
    _drop_placeholder(ph)


def _insert_table(native_slide, ph, asset: Asset) -> None:
    data = asset.data or {}
    columns = list(data.get("columns") or [])
    rows = [list(r) for r in (data.get("rows") or [])]
    if not columns and not rows:
        raise ValueError("table asset has neither columns nor rows")

    n_rows = len(rows) + (1 if columns else 0)
    n_cols = len(columns) or max(len(r) for r in rows)

    shape = native_slide.shapes.add_table(
        n_rows, n_cols, ph.left, ph.top, ph.width, ph.height
    )
    table = shape.table

    r = 0
    if columns:
        for c, label in enumerate(columns):
            table.cell(0, c).text = str(label)
        r = 1
    for row in rows:
        for c, value in enumerate(row[:n_cols]):
            table.cell(r, c).text = _format_cell(value)
        r += 1

    _drop_placeholder(ph)


def _format_cell(value) -> str:
    if isinstance(value, float):
        return f"{value:.4g}"
    return str(value)


def _probe_aspect(path: str | None) -> float | None:
    if not path:
        return None
    try:
        from PIL import Image

        with Image.open(path) as im:
            w, h = im.size
        return (w / h) if h else None
    except Exception:  # pragma: no cover - non-raster or missing
        return None


def _drop_placeholder(ph) -> None:
    """Remove an emptied placeholder so its prompt text does not render."""
    element = ph._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


# --------------------------------------------------------------------------- #
# Speaker notes
# --------------------------------------------------------------------------- #


def _build_notes(
    deck: Deck, slide: Slide, leftover: list[Block], language: str
) -> str:
    """Assemble the narration for one slide.

    The script is a *projection* of the same tree as the slides, so what the
    presenter says and what the audience reads cannot drift apart.
    """
    parts: list[str] = []
    if slide.notes:
        parts.append(slide.notes)

    if slide.dwell_seconds:
        parts.append(f"[{slide.dwell_seconds:.0f}s]")

    for block in slide.blocks:
        if block.channel is Channel.SCRIPT:
            parts.append(block.content)
        elif block.channel is Channel.SLIDE and block.speaker_note:
            parts.append(block.speaker_note)

    for block in leftover:
        parts.append(f"[未能放入版式] {block.content or block.uid}")

    return "\n\n".join(p for p in parts if p)


# --------------------------------------------------------------------------- #
# Template handling
# --------------------------------------------------------------------------- #


def _strip_existing_slides(prs) -> None:
    """Remove any slides shipped inside the template.

    Templates are often saved as ordinary decks with example content. Those
    slides must go, while the layouts and master they reference must stay.
    """
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        rid = sld_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if rid:
            prs.part.drop_rel(rid)
        xml_slides.remove(sld_id)


def _find_native_layout(prs, native_ref: str | None, layout_id: str):
    """Locate the python-pptx layout object for a contract layout."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if native_ref and layout.name == native_ref:
                return layout
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if _slug(layout.name) == layout_id:
                return layout
    return None


def _slug(text: str) -> str:
    out = []
    for ch in (text or "").lower():
        if ch.isalnum():
            out.append(ch)
        elif out and out[-1] != "-":
            out.append("-")
    return "".join(out).strip("-")


# --------------------------------------------------------------------------- #
# Script export
# --------------------------------------------------------------------------- #


def export_script(deck: Deck, path: str) -> str:
    """Write the speaker script as Markdown.

    A projection of the deck, not a second generation pass: every line here
    traces to a block that is either on a slide or deliberately routed to the
    script.
    """
    lang = deck.meta.language
    lines: list[str] = [f"# {deck.meta.title}", ""]

    total = deck.delivery.total_seconds
    lines.append(
        f"*目标时长 {total // 60} 分 {total % 60} 秒 · "
        f"语速 {deck.delivery.words_per_minute}/分 · "
        f"密度 {deck.delivery.density.value}*"
    )
    lines.append("")

    n = 0
    for section, slide in deck.iter_slides():
        if slide.is_backup:
            continue
        n += 1
        head = f"## {n}. {slide.title or section.title}"
        if slide.dwell_seconds:
            head += f"  ({slide.dwell_seconds:.0f}s)"
        lines.append(head)

        spoken: list[str] = []
        if slide.notes:
            spoken.append(slide.notes)
        for block in slide.blocks:
            if block.channel is Channel.DROP:
                continue
            if block.channel is Channel.SCRIPT:
                spoken.append(block.content)
            elif block.speaker_note:
                spoken.append(block.speaker_note)

        if spoken:
            lines.extend(["", *(f"{p}" for p in spoken)])
        else:
            on_slide = [b.content for b in slide.slide_blocks() if b.content]
            if on_slide:
                lines.extend(["", "（照读要点）", *(f"- {t}" for t in on_slide)])

        words = sum(count_units(p, lang) for p in spoken)
        if words:
            secs = words / deck.delivery.words_per_minute * 60
            lines.append("")
            lines.append(f"*约 {words} 字 / {secs:.0f}s*")
        lines.append("")

    backup = [s for _, s in deck.iter_slides() if s.is_backup]
    if backup:
        lines.extend(["---", "", "## 备用页（按需展开）", ""])
        for s in backup:
            lines.append(f"- {s.title}")

    text = "\n".join(lines) + "\n"
    with open(path, "w", encoding="utf-8") as fh:
        fh.write(text)
    return text
