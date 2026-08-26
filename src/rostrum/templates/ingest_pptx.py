"""PPTX template ingestion: layouts -> measured capacity contract.

What this does and why
----------------------
A template is not a colour scheme. To pour arbitrary content into an arbitrary
user-supplied template, the renderer needs to know, for every layout, *which
functional role it can serve* and *how much content each of its slots actually
holds*. Neither is stated in the file; both must be derived.

Two derivations happen here:

1. **Role classification.** Placeholder composition (one title + one body, title
   + two objects, picture-bearing, ...) maps onto :class:`SlideRole`. This is
   what lets the renderer match role -> layout instead of index -> index, and
   therefore what makes "upload your own lab template" work at all.

2. **Capacity measurement.** Slot geometry is resolved through the
   placeholder-inheritance chain (layout -> master), then handed to
   :mod:`rostrum.measure.text` for real glyph measurement. The resulting
   ``capacity_units`` is a *sound upper bound*, which is what allows the planner
   to prevent overflow rather than detect it afterwards.
"""

from __future__ import annotations

import contextlib
from dataclasses import dataclass, field

from lxml import etree

from rostrum.ir.enums import Renderer, SlideRole
from rostrum.measure.text import (
    LINE_HEIGHT_CJK,
    capacity_units,
    emu_to_pt,
    lines_available,
    load_font,
)
from rostrum.templates.contract import Box, Layout, Slot, TemplateContract

try:
    from pptx import Presentation

    _HAVE_PPTX = True
except ImportError:  # pragma: no cover - optional dependency
    _HAVE_PPTX = False


# Placeholder types that carry running text.
_TEXT_PH = {
    "TITLE",
    "CENTER_TITLE",
    "SUBTITLE",
    "BODY",
    "OBJECT",
    "VERTICAL_BODY",
    "VERTICAL_TITLE",
}
_TITLE_PH = {"TITLE", "CENTER_TITLE", "VERTICAL_TITLE"}
_VISUAL_PH = {"PICTURE", "CHART", "TABLE", "MEDIA_CLIP", "OBJECT"}
_CHROME_PH = {"DATE", "FOOTER", "SLIDE_NUMBER"}

# Nominal font sizes (pt) when a template declares none. PowerPoint's own
# defaults for the corresponding placeholder levels.
_DEFAULT_SIZE = {"title": 44.0, "subtitle": 32.0, "body": 18.0, "footer": 12.0}

# Legibility floors for academic projection. Shrinking past these to force a fit
# is a defect, not a solution: the back row cannot read 12pt body text.
_MIN_SIZE = {"title": 24.0, "subtitle": 18.0, "body": 14.0, "footer": 8.0}


@dataclass
class IngestReport:
    """Diagnostics from ingesting a template."""

    template_id: str
    layouts_found: int
    layouts_kept: int
    roles_covered: set[SlideRole]
    unmeasured_slots: int = 0
    warnings: list[str] = field(default_factory=list)

    def missing_for(self, needed: set[SlideRole]) -> set[SlideRole]:
        return needed - self.roles_covered


def ingest_pptx(
    path: str,
    *,
    template_id: str | None = None,
    name: str | None = None,
    language: str = "zh",
    font_path: str | None = None,
    license: str | None = None,
) -> tuple[TemplateContract, IngestReport]:
    """Parse a ``.pptx`` template into a measured :class:`TemplateContract`.

    Parameters
    ----------
    path:
        Template file. Only its layouts and master are read; any existing
        slides are ignored.
    font_path:
        Font used for measurement. Should match what the deck will render with;
        a substituted font invalidates the measurements.
    """
    if not _HAVE_PPTX:  # pragma: no cover
        raise RuntimeError("python-pptx is required: pip install 'rostrum[pptx]'")

    prs = Presentation(path)
    tid = template_id or _slugify(name or _basename(path))
    font = load_font(font_path)

    page_w_pt = emu_to_pt(prs.slide_width)
    page_h_pt = emu_to_pt(prs.slide_height)

    report = IngestReport(
        template_id=tid,
        layouts_found=len(prs.slide_layouts),
        layouts_kept=0,
        roles_covered=set(),
    )
    if not font.measured and font_path:
        report.warnings.append(
            f"font {font_path!r} could not be parsed; capacities fall back to "
            "nominal advance widths and are less accurate"
        )
    elif font_path is None:
        report.warnings.append(
            "no font supplied; capacities use nominal advance widths. Pass "
            "font_path matching the render font for exact bounds"
        )

    layouts: list[Layout] = []
    for index, native in enumerate(prs.slide_layouts):
        layout = _ingest_layout(
            native,
            index=index,
            page_w_pt=page_w_pt,
            page_h_pt=page_h_pt,
            language=language,
            font=font,
            report=report,
        )
        if layout is not None:
            layouts.append(layout)
            report.roles_covered |= set(layout.roles)

    if not layouts:
        raise ValueError(
            f"{path}: no usable layouts found. A template needs at least one "
            "layout with a title or body placeholder"
        )

    report.layouts_kept = len(layouts)
    contract = TemplateContract(
        template_id=tid,
        name=name or _basename(path),
        renderer=Renderer.PPTX,
        source_path=path,
        page_aspect=_aspect_label(page_w_pt, page_h_pt),
        layouts=layouts,
        license=license,
        fonts=sorted(_fonts_used(prs)),
    )
    return contract, report


# --------------------------------------------------------------------------- #
# Layout ingestion
# --------------------------------------------------------------------------- #


def _ingest_layout(
    native,
    *,
    index: int,
    page_w_pt: float,
    page_h_pt: float,
    language: str,
    font,
    report: IngestReport,
) -> Layout | None:
    slots: list[Slot] = []
    kinds: list[str] = []

    for ph in native.placeholders:
        ph_type = _ph_type_name(ph)
        if ph_type in _CHROME_PH:
            continue  # date/footer/page-number are template chrome, not content

        geom = _resolve_geometry(ph, native)
        if geom is None:
            report.unmeasured_slots += 1
            report.warnings.append(
                f"layout {native.name!r}: placeholder "
                f"{ph.placeholder_format.idx} has no resolvable geometry; skipped"
            )
            continue

        x_pt, y_pt, w_pt, h_pt = geom
        kind = _slot_kind(ph_type, ph)
        kinds.append(kind)

        size = _nominal_size(ph, kind, height_pt=h_pt)
        slot = Slot(
            slot_id=f"ph{ph.placeholder_format.idx}",
            kind=kind,
            box=Box(
                x=_clamp(x_pt / page_w_pt),
                y=_clamp(y_pt / page_h_pt),
                w=_clamp(w_pt / page_w_pt, lo=0.01),
                h=_clamp(h_pt / page_h_pt, lo=0.01),
            ),
            font_size_pt=size,
            min_font_size_pt=_MIN_SIZE.get(kind),
            required=kind == "title",
        )

        if kind in {"title", "subtitle", "body", "footer"}:
            slot.capacity_units = capacity_units(
                width_pt=w_pt,
                height_pt=h_pt,
                font_size_pt=size,
                font=font,
                language=language,
                bullet_indent_pt=_BULLET_INDENT_PT if kind == "body" else 0.0,
            )
            slot.capacity_lines = lines_available(
                height_pt=h_pt, font_size_pt=size
            )
            slot.max_bullet_level = min(3, max(1, (slot.capacity_lines or 1) // 2))

        slots.append(slot)

    if not slots:
        return None

    roles = _classify(native.name, kinds, slots)
    if not roles:
        return None

    return Layout(
        layout_id=_slugify(native.name) or f"layout-{index}",
        roles=roles,
        slots=slots,
        native_ref=native.name,
    )


# Body placeholders carry a bullet glyph and hanging indent that eat horizontal
# space. PowerPoint's default first-level indent.
_BULLET_INDENT_PT = 18.0


def _resolve_geometry(ph, layout) -> tuple[float, float, float, float] | None:
    """Resolve a placeholder's box, following the inheritance chain.

    Layout placeholders frequently declare no geometry of their own and inherit
    it from the slide master. Reading ``ph.left`` naively yields ``None`` and
    silently produces a zero-sized slot -- which would make every capacity zero
    and every slide overflow.
    """
    box = _explicit_box(ph)
    if box is not None:
        return box

    idx = ph.placeholder_format.idx
    master = getattr(layout, "slide_master", None)
    if master is not None:
        with contextlib.suppress(KeyError):
            for mph in master.placeholders:
                if mph.placeholder_format.idx == idx:
                    box = _explicit_box(mph)
                    if box is not None:
                        return box
        # Body geometry is commonly inherited from the master's body
        # placeholder even when indices differ.
        with contextlib.suppress(Exception):
            for mph in master.placeholders:
                if _ph_type_name(mph) in {"BODY", "OBJECT"}:
                    box = _explicit_box(mph)
                    if box is not None:
                        return box
    return None


def _explicit_box(ph) -> tuple[float, float, float, float] | None:
    try:
        left, top, width, height = ph.left, ph.top, ph.width, ph.height
    except (AttributeError, KeyError):  # pragma: no cover
        return None
    if None in (left, top, width, height):
        return None
    if width <= 0 or height <= 0:
        return None
    return (
        emu_to_pt(left),
        emu_to_pt(top),
        emu_to_pt(width),
        emu_to_pt(height),
    )


def _ph_type_name(ph) -> str:
    try:
        t = ph.placeholder_format.type
    except (AttributeError, KeyError):  # pragma: no cover
        return "BODY"
    if t is None:
        return "BODY"
    name = getattr(t, "name", None)
    return name or str(t).split()[0].strip("<>")


def _slot_kind(ph_type: str, ph) -> str:
    if ph_type in _TITLE_PH:
        return "title"
    if ph_type == "SUBTITLE":
        return "subtitle"
    if ph_type == "PICTURE":
        return "figure"
    if ph_type == "TABLE":
        return "table"
    if ph_type == "CHART":
        return "figure"
    # OBJECT placeholders accept either text or a picture. Treat them as body
    # text, since that is the dominant academic use, while _classify still lets
    # the layout serve figure roles.
    return "body"


def _nominal_size(ph, kind: str, height_pt: float | None = None) -> float:
    """Font size for measurement: declared by the template, else derived.

    When a template declares no size we cannot simply apply PowerPoint's default
    (44pt for a title): a deliberately short title bar would then measure as
    holding zero lines, and the planner would reject every title as
    unfittable. Deriving a size that yields at least one line from the actual box
    height keeps the capacity meaningful, floored at the legibility minimum.
    """
    with contextlib.suppress(AttributeError, KeyError, IndexError):
        tf = ph.text_frame
        for para in tf.paragraphs:
            if para.font.size is not None:
                return float(para.font.size.pt)
            for run in para.runs:
                if run.font.size is not None:
                    return float(run.font.size.pt)

    default = _DEFAULT_SIZE.get(kind, 18.0)
    if height_pt is None or height_pt <= 0:
        return default

    # Largest size at which one line still fits the box.
    fits = height_pt / LINE_HEIGHT_CJK
    if fits >= default:
        return default
    return max(_MIN_SIZE.get(kind, 12.0), round(fits, 1))


# --------------------------------------------------------------------------- #
# Role classification
# --------------------------------------------------------------------------- #

# Name hints take precedence over composition: a template author who called a
# layout "Section Header" knows better than our heuristics.
_NAME_HINTS: list[tuple[tuple[str, ...], list[SlideRole]]] = [
    (("title slide", "cover", "封面", "标题幻灯片"), [SlideRole.COVER]),
    (("section", "章节", "节标题"), [SlideRole.SECTION]),
    (("agenda", "outline", "contents", "目录", "大纲"), [SlideRole.AGENDA]),
    (("comparison", "对比", "两栏", "two content"), [SlideRole.TWO_COLUMN]),
    (("three", "三栏"), [SlideRole.THREE_COLUMN]),
    (("picture", "image", "figure", "图片", "图"), [SlideRole.BIG_FIGURE]),
    (("table", "表格"), [SlideRole.TABLE]),
    (("equation", "formula", "公式"), [SlideRole.EQUATION]),
    (("timeline", "时间线", "计划"), [SlideRole.TIMELINE]),
    (("summary", "conclusion", "总结", "结论"), [SlideRole.SUMMARY]),
    (("thank", "acknowledg", "致谢", "感谢"), [SlideRole.ACKNOWLEDGEMENT]),
    (("backup", "reserve", "备用", "附录", "appendix"), [SlideRole.BACKUP]),
    (("blank", "空白"), []),
]


def _classify(name: str, kinds: list[str], slots: list[Slot]) -> list[SlideRole]:
    """Infer which functional roles a layout can serve.

    A layout usually serves several: a title-plus-body layout is a fine home for
    a text-dense slide, a summary, or an agenda. Being generous here is what
    keeps a sparse user template usable, while ``missing_roles`` still reports
    genuine gaps before rendering.
    """
    lowered = (name or "").lower()
    for needles, roles in _NAME_HINTS:
        if any(n in lowered for n in needles):
            if not roles:
                return []  # explicitly blank layout
            return _augment(roles, kinds, slots)

    n_title = kinds.count("title")
    n_body = kinds.count("body")
    n_figure = kinds.count("figure")
    n_sub = kinds.count("subtitle")

    if n_title and n_sub and not n_body:
        return [SlideRole.COVER]
    if n_title and not n_body and not n_figure:
        return [SlideRole.SECTION, SlideRole.ACKNOWLEDGEMENT]
    if n_body >= 3:
        return [SlideRole.THREE_COLUMN, SlideRole.TEXT_DENSE]
    if n_body == 2:
        return [SlideRole.TWO_COLUMN, SlideRole.TEXT_FIGURE, SlideRole.TEXT_DENSE]
    if n_figure and n_body:
        return [SlideRole.TEXT_FIGURE, SlideRole.BIG_FIGURE, SlideRole.TABLE]
    if n_figure and not n_body:
        return [SlideRole.BIG_FIGURE, SlideRole.TABLE]
    if n_body == 1:
        return _augment(
            [
                SlideRole.TEXT_DENSE,
                SlideRole.AGENDA,
                SlideRole.SUMMARY,
                SlideRole.EQUATION,
            ],
            kinds,
            slots,
        )
    return []


def _augment(
    roles: list[SlideRole], kinds: list[str], slots: list[Slot]
) -> list[SlideRole]:
    """Add roles a layout can additionally serve, given its composition."""
    out = list(roles)
    body_slots = [s for s in slots if s.kind == "body"]

    # A single generous body slot can host a figure: an OBJECT placeholder in
    # PowerPoint accepts a picture as readily as text.
    if len(body_slots) == 1 and (body_slots[0].box.w * body_slots[0].box.h) > 0.18:
        for extra in (SlideRole.TEXT_FIGURE, SlideRole.TABLE):
            if extra not in out:
                out.append(extra)
    if len(body_slots) >= 2 and SlideRole.TWO_COLUMN not in out:
        out.append(SlideRole.TWO_COLUMN)
    return out


# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #


def _fonts_used(prs) -> set[str]:
    """Fonts a template will actually render with, for substitution warnings.

    Two sources, both needed. Shapes may name a face directly, but a
    well-formed template instead references the theme (``+mj-lt`` / ``+mn-lt``)
    and declares the real faces once in the theme part -- which is the correct
    way to do it, and invisible to a scan of shape properties alone. Reading only
    shapes reported no fonts at all for exactly such a template.
    """
    found: set[str] = set()
    with contextlib.suppress(Exception):
        for master in prs.slide_masters:
            for ph in master.placeholders:
                for para in ph.text_frame.paragraphs:
                    if para.font.name and not para.font.name.startswith("+"):
                        found.add(para.font.name)
                    for run in para.runs:
                        if run.font.name and not run.font.name.startswith("+"):
                            found.add(run.font.name)
    found |= _theme_fonts(prs)
    return found


def _theme_fonts(prs) -> set[str]:
    """Faces declared in the theme's font scheme, including the CJK entries.

    ``script="Hans"`` is what governs Chinese glyph selection, so it matters more
    than the latin typeface for the decks this tool produces.
    """
    found: set[str] = set()
    with contextlib.suppress(Exception):
        part = prs.slide_masters[0].part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        root = etree.fromstring(part.blob)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        for scheme in ("majorFont", "minorFont"):
            node = root.find(f".//a:fontScheme/a:{scheme}", ns)
            if node is None:
                continue
            for tag in ("latin", "ea"):
                el = node.find(f"a:{tag}", ns)
                if el is not None and el.get("typeface"):
                    found.add(el.get("typeface"))
            for font in node.findall("a:font", ns):
                if font.get("script") in ("Hans", "Hant") and font.get("typeface"):
                    found.add(font.get("typeface"))
    return {f for f in found if f and not f.startswith("+")}


def _aspect_label(w_pt: float, h_pt: float) -> str:
    if h_pt <= 0:  # pragma: no cover
        return "16:9"
    r = w_pt / h_pt
    for label, value in (("16:9", 16 / 9), ("16:10", 1.6), ("4:3", 4 / 3), ("3:2", 1.5)):
        if abs(r - value) < 0.02:
            return label
    return f"{r:.3f}:1"


def _clamp(v: float, lo: float = 0.0, hi: float = 1.0) -> float:
    return max(lo, min(hi, v))


def _slugify(text: str) -> str:
    out = []
    for ch in (text or "").lower():
        if ch.isalnum():
            out.append(ch)
        elif out and out[-1] != "-":
            out.append("-")
    return "".join(out).strip("-")


def _basename(path: str) -> str:
    return path.replace("\\", "/").rsplit("/", 1)[-1].rsplit(".", 1)[0]
